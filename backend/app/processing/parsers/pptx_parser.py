import io

from pptx import Presentation

from app.core.exceptions import ValidationException
from app.processing.parsers.base_parser import BaseParser


class PptxParser(BaseParser):
    """Extracts text content from PowerPoint (.pptx) files."""

    def extract_text(self, file_bytes: bytes) -> str:
        try:
            presentation = Presentation(io.BytesIO(file_bytes))
        except Exception as exc:
            raise ValidationException("Could not read PPTX file — it may be corrupted") from exc

        slides_text = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_lines = [f"--- Slide {slide_number} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs)
                        if text.strip():
                            slide_lines.append(text)
            if len(slide_lines) > 1:
                slides_text.append("\n".join(slide_lines))

        return "\n\n".join(slides_text).strip()