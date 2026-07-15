"""Generates sample fixture files for each supported document type.

Run this once manually: python tests/generate_fixtures.py
Output files land in ../sample-data/
"""
import csv
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

SAMPLE_DIR = Path(__file__).resolve().parent.parent.parent / "sample-data"
SAMPLE_DIR.mkdir(exist_ok=True)


def generate_txt() -> None:
    (SAMPLE_DIR / "sample.txt").write_text(
        "This is a sample text document.\nIt has multiple lines.\nUsed for testing the TXT parser.",
        encoding="utf-8",
    )


def generate_csv() -> None:
    with open(SAMPLE_DIR / "sample.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["name", "department", "salary"])
        writer.writerow(["Alice", "Engineering", "95000"])
        writer.writerow(["Bob", "Sales", "72000"])


def generate_docx() -> None:
    doc = DocxDocument()
    doc.add_heading("Sample Company Policy", level=1)
    doc.add_paragraph("This document outlines the sample remote work policy.")
    doc.add_paragraph("Employees may work remotely up to three days per week.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Policy"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "Remote days"
    table.rows[1].cells[1].text = "3"
    doc.save(SAMPLE_DIR / "sample.docx")


def generate_pptx() -> None:
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Q3 Company Overview"
    body = slide.placeholders[1]
    body.text_frame.text = "Revenue grew 15% year over year."
    presentation.save(SAMPLE_DIR / "sample.pptx")


def generate_xlsx() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Budget"
    sheet.append(["Category", "Q1", "Q2"])
    sheet.append(["Marketing", 10000, 12000])
    sheet.append(["Engineering", 50000, 55000])
    workbook.save(SAMPLE_DIR / "sample.xlsx")


def generate_eml() -> None:
    content = (
        "From: manager@company.com\n"
        "To: team@company.com\n"
        "Subject: Weekly Update\n\n"
        "Hi team,\n\nHere is this week's project update. Everything is on track.\n"
    )
    (SAMPLE_DIR / "sample.eml").write_text(content, encoding="utf-8")


if __name__ == "__main__":
    generate_txt()
    generate_csv()
    generate_docx()
    generate_pptx()
    generate_xlsx()
    generate_eml()
    print(f"Sample fixtures generated in: {SAMPLE_DIR}")
    